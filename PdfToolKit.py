import tkinter as tk
from PIL import Image, ImageTk
from tkinter import filedialog, messagebox, PhotoImage
from pdf2docx import Converter 
from PyPDF2 import PdfMerger, PdfReader, PdfWriter, PdfFileWriter, PdfFileReader
import fitz
from pptx import Presentation
from pptx.util import Inches
import pytesseract

def toggle_mode():
    global is_day_mode
    if is_day_mode:
        # Cambiar a modo noche
        frame.config(bg="black")
        for button in buttons:
            button.config(bg="gray", fg="white")
        toggle_button.config(image=day_icon, bg="black", fg="white", text="Modo Día")
    else:
        # Cambiar a modo día
        frame.config(bg="white")
        for button in buttons:
            button.config(bg="lightgray", fg="black")
        toggle_button.config(image=night_icon, bg="white", fg="black", text="Modo Noche")
    
    is_day_mode = not is_day_mode

def resize_image(image_path, size):
    image = Image.open(image_path).convert("RGBA")
    datas = image.getdata()
    new_data = []
    for item in datas:
        if item[0] > 200 and item[1] > 200 and item[2] > 200:  # Threshold for white
            new_data.append((255, 255, 255, 0))
        else:
            new_data.append(item)
    image.putdata(new_data)
    image = image.resize(size, Image.Resampling.LANCZOS)
    return ImageTk.PhotoImage(image)

root = tk.Tk()
root.title("PDF-Tool-Kit")

# Configuración inicial del modo
is_day_mode = True

# Crear un Canvas (aunque no se está usando en el código, podría ser útil)
canvas = tk.Canvas(root, width=20, height=20)
canvas.pack()

# Frame para organizar los botones en una cuadrícula
frame = tk.Frame(root, bg="white")
frame.pack(fill=tk.BOTH, expand=True)

# Label para mostrar los mensajes
message_label = tk.Label(root, text="")
message_label.pack(pady=5)

# Variable para almacenar la ruta del archivo PDF
pdf_path_var = tk.StringVar()

# Cargar los iconos
icons = [f"img/image{i}.png" for i in range(1, 26)]
images = [resize_image(icon, (60, 60)) for icon in icons]

day_icon = resize_image("img/day_icon.png", (30, 30))
night_icon = resize_image("img/night_icon.png", (30, 30))

# Crear botones
buttons = []

def create_button(text, image, command):
    button = tk.Button(frame, text=text, image=image, compound="top", command=command, highlightthickness=1)
    return button

# Función para cada opción
def opcion1():
    message_label.config(text="Convertir PDF a Word seleccionada")
    
    def select_pdf():
        file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file_path:
            pdf_path_var.set(file_path)

    def convert_pdf_to_word():
        pdf_path = pdf_path_var.get()
        if not pdf_path:
            messagebox.showerror("Error", "Por favor, seleccione un archivo PDF.")
            return
        
        save_path = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word files", "*.docx")])
        if save_path:
            try:
                cv = Converter(pdf_path)
                cv.convert(save_path, start=0, end=None)
                cv.close()
                messagebox.showinfo("Éxito", f"Archivo convertido y guardado en: {save_path}")
            except Exception as e:
                messagebox.showerror("Error", f"No se pudo convertir el archivo: {e}")

    pdf_window = tk.Toplevel(root)
    pdf_window.title("Convertir PDF a Word")

    tk.Label(pdf_window, text="Archivo PDF:").grid(row=0, column=0, padx=10, pady=10)
    tk.Entry(pdf_window, textvariable=pdf_path_var, width=50).grid(row=0, column=1, padx=10, pady=10)

    tk.Button(pdf_window, text="Seleccionar PDF", command=select_pdf).grid(row=0, column=2, padx=10, pady=10)
    tk.Button(pdf_window, text="Convertir a Word", command=convert_pdf_to_word).grid(row=1, columnspan=3, pady=20)

def opcion2():
    message_label.config(text="Unir PDF seleccionada")
    
    pdf_files = []

    def select_pdfs():
        files = filedialog.askopenfilenames(filetypes=[("PDF files", "*.pdf")])
        if files:
            pdf_files.extend(files)
            message_label.config(text=f"{len(pdf_files)} archivos seleccionados")

    def merge_pdfs():
        if not pdf_files:
            messagebox.showerror("Error", "Por favor, seleccione al menos dos archivos PDF.")
            return

        save_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if save_path:
            try:
                merger = PdfMerger()
                for pdf in pdf_files:
                    merger.append(pdf)
                merger.write(save_path)
                merger.close()
                messagebox.showinfo("Éxito", f"Archivos unidos y guardados en: {save_path}")
                pdf_files.clear()
                message_label.config(text="")
            except Exception as e:
                messagebox.showerror("Error", f"No se pudo unir los archivos: {e}")

    merge_window = tk.Toplevel(root)
    merge_window.title("Unir PDF")

    tk.Button(merge_window, text="Seleccionar PDFs", command=select_pdfs).pack(pady=10)
    tk.Button(merge_window, text="Unir PDFs", command=merge_pdfs).pack(pady=10)

def opcion3():
    message_label.config(text="Dividir PDF seleccionada")
    
    def select_pdf():
        file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file_path:
            pdf_path_var.set(file_path)

    def split_pdf():
        pdf_path = pdf_path_var.get()
        if not pdf_path:
            messagebox.showerror("Error", "Por favor, seleccione un archivo PDF.")
            return
        
        try:
            reader = PdfReader(pdf_path)
            num_pages = len(reader.pages)
            
            for i in range(num_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                
                save_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")], initialfile=f"page_{i + 1}.pdf")
                if save_path:
                    with open(save_path, 'wb') as output_pdf:
                        writer.write(output_pdf)
            
            messagebox.showinfo("Éxito", "Archivo dividido y guardado.")
        except Exception as e:
            messagebox.showerror("Error", f"No se pudo dividir el archivo: {e}")

    split_window = tk.Toplevel(root)
    split_window.title("Dividir PDF")

    tk.Label(split_window, text="Archivo PDF:").grid(row=0, column=0, padx=10, pady=10)
    tk.Entry(split_window, textvariable=pdf_path_var, width=50).grid(row=0, column=1, padx=10, pady=10)

    tk.Button(split_window, text="Seleccionar PDF", command=select_pdf).grid(row=0, column=2, padx=10, pady=10)
    tk.Button(split_window, text="Dividir PDF", command=split_pdf).grid(row=1, columnspan=3, pady=20)

def opcion4():
    message_label.config(text="Convertir PDF a PPT seleccionada")

    def select_pdf():
        file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file_path:
            pdf_path_var.set(file_path)

    def convert_pdf_to_ppt():
        pdf_path = pdf_path_var.get()
        if not pdf_path:
            messagebox.showerror("Error", "Por favor, seleccione un archivo PDF.")
            return
        
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        if save_path:
            try:
                pdf_document = fitz.open(pdf_path)
                presentation = Presentation()

                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    pix = page.get_pixmap()
                    image_path = f"page_{page_num + 1}.png"
                    pix.save(image_path)
                    
                    slide_layout = presentation.slide_layouts[5]
                    slide = presentation.slides.add_slide(slide_layout)
                    left = top = Inches(0)
                    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(10), height=Inches(7.5))

                presentation.save(save_path)
                messagebox.showinfo("Éxito", f"Archivo convertido y guardado en: {save_path}")
            except Exception as e:
                messagebox.showerror("Error", f"No se pudo convertir el archivo: {e}")

    pdf_window = tk.Toplevel(root)
    pdf_window.title("Convertir PDF a PPT")

    tk.Label(pdf_window, text="Archivo PDF:").grid(row=0, column=0, padx=10, pady=10)
    tk.Entry(pdf_window, textvariable=pdf_path_var, width=50).grid(row=0, column=1, padx=10, pady=10)

    tk.Button(pdf_window, text="Seleccionar PDF", command=select_pdf).grid(row=0, column=2, padx=10, pady=10)
    tk.Button(pdf_window, text="Convertir a PPT", command=convert_pdf_to_ppt).grid(row=1, columnspan=3, pady=20)

def opcion5():
    message_label.config(text="PDF a PowerPoint")
    def select_pdf():
        file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file_path:
            pdf_path_var.set(file_path)

    def convert_pdf_to_ppt():
        pdf_path = pdf_path_var.get()
        if not pdf_path:
            messagebox.showerror("Error", "Por favor, seleccione un archivo PDF.")
            return
        
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        if save_path:
            try:
                pdf_document = fitz.open(pdf_path)
                presentation = Presentation()

                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    pix = page.get_pixmap()
                    image_path = f"temp_page_{page_num + 1}.png"
                    pix.save(image_path)

                    slide_layout = presentation.slide_layouts[5]
                    slide = presentation.slides.add_slide(slide_layout)

                    title_placeholder = slide.shapes.title
                    title_placeholder.text = f"Page {page_num + 1}"

                    left = Inches(1)
                    top = Inches(1.5)
                    slide.shapes.add_picture(image_path, left, top, width=Inches(8.5), height=Inches(9))

                presentation.save(save_path)
                messagebox.showinfo("Éxito", f"Archivo convertido y guardado en: {save_path}")
            except Exception as e:
                messagebox.showerror("Error", f"No se pudo convertir el archivo: {e}")

    # Crear una nueva ventana para la selección y conversión de PDF a PowerPoint
    ppt_window = tk.Toplevel(root)
    ppt_window.title("Convertir PDF a PowerPoint")

    tk.Label(ppt_window, text="Archivo PDF:").grid(row=0, column=0, padx=10, pady=10)
    tk.Entry(ppt_window, textvariable=pdf_path_var, width=50).grid(row=0, column=1, padx=10, pady=10)

    # Botón para seleccionar el archivo PDF
    tk.Button(ppt_window, text="Seleccionar PDF", command=select_pdf).grid(row=0, column=2, padx=10, pady=10)

    # Botón para iniciar la conversión
    tk.Button(ppt_window, text="Convertir a PowerPoint", command=convert_pdf_to_ppt).grid(row=1, columnspan=3, pady=20)
    pass

def opcion6():
    message_label.config(text="PDF a Excel")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion7():
    message_label.config(text="Word a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion8():
    message_label.config(text="PowerPoint a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion9():
    message_label.config(text="Excel a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion10():
    message_label.config(text="Editar PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion11():
    message_label.config(text="PDF a JPG")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion12():
    message_label.config(text="JPG a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion13(): # ME DA ERROR --> PENDIENTE  DE REVISAR
    message_label.config(text="Firmar PDF")
    # Aquí iría la lógica específica para esta opción

    file_path = filedialog.askopenfilename(title="Selecciona un PDF", filetypes=[("PDF files", "*.pdf")])
    if not file_path:
        return

    save_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
    if not save_path:
        return

    try:
        # Aquí deberías agregar la lógica para firmar el PDF
        pdf_reader = PdfFileReader(file_path)
        pdf_writer = PdfFileWriter()

        # Copiar todas las páginas del PDF original al nuevo PDF
        for page_num in range(pdf_reader.getNumPages()):
            page = pdf_reader.getPage(page_num)
            pdf_writer.addPage(page)

        # Simular una firma añadiendo una marca de agua o similar
        # Aquí puedes añadir tu propia lógica para firmar el PDF

        with open(save_path, "wb") as output_pdf:
            pdf_writer.write(output_pdf)

        messagebox.showinfo("Éxito", "El PDF ha sido firmado y guardado exitosamente.")
    except Exception as e:
        messagebox.showerror("Error", f"Ha ocurrido un error al firmar el PDF: {str(e)}")
    pass

def opcion14():
    message_label.config(text="Marca de agua")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion15():
    message_label.config(text="Rotar PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion16():
    message_label.config(text="Html a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion17():
    message_label.config(text="Desbloquear PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion18():
    message_label.config(text="Proteger PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion19():
    message_label.config(text="Ordenar PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion20():
    message_label.config(text="PDF a PDF/a")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion21():
    message_label.config(text="Repara PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion22():
    message_label.config(text="Enumerar páginas")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion23():
    message_label.config(text="Escanea a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

def opcion24(): # ME DA ERROR --> PENDIENTE  DE REVISAR
    message_label.config(text="OCR PDF")
    
    file_path = filedialog.askopenfilename(title="Selecciona un PDF", filetypes=[("PDF files", "*.pdf")])
    if not file_path:
        return

    try:
        pdf_reader = PdfFileReader(file_path)
        extracted_text = ""

        for page_num in range(pdf_reader.getNumPages()):
            page = pdf_reader.getPage(page_num)
            xObject = page['/Resources']['/XObject'].getObject()
            for obj in xObject:
                if xObject[obj]['/Subtype'] == '/Image':
                    size = (xObject[obj]['/Width'], xObject[obj]['/Height'])
                    data = xObject[obj]._data
                    mode = "RGB" if xObject[obj]['/ColorSpace'] == '/DeviceRGB' else "P"

                    img = Image.frombytes(mode, size, data)
                    text = pytesseract.image_to_string(img)
                    extracted_text += text + "\n"

        if extracted_text:
            save_path = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("Text files", "*.txt")])
            if save_path:
                with open(save_path, "w", encoding="utf-8") as text_file:
                    text_file.write(extracted_text)
                messagebox.showinfo("Éxito", "El texto ha sido extraído y guardado exitosamente.")
        else:
            messagebox.showinfo("Resultado", "No se ha encontrado texto en el PDF.")

    except Exception as e:
        messagebox.showerror("Error", f"Ha ocurrido un error al realizar OCR en el PDF: {str(e)}")
    pass

def opcion25():
    message_label.config(text="Comparar a PDF")
    # Aquí iría la lógica específica para esta opción
    pass

# Lista de opciones con texto, imagen y comando
opciones = [
    ("Convertir PDF a Word", images[0], opcion1),
    ("Unir PDF", images[1], opcion2),
    ("Dividir PDF", images[2], opcion3),
    ("Comprimir PDF", images[3], opcion4),
    ("PDF a PowerPoint", images[4], opcion5),
    ("PDF a Excel", images[5], opcion6),
    ("Word a PDF", images[6], opcion7),
    ("PowerPoint a PDF", images[7], opcion8),
    ("Excel a PDF", images[8], opcion9),
    ("Editar PDF", images[9], opcion10),
    ("PDF a JPG", images[10], opcion11),
    ("JPG a PDF", images[11], opcion12),
    ("Firmar PDF", images[12], opcion13),
    ("Marca de agua", images[13], opcion14),
    ("Rotar PDF", images[14], opcion15),
    ("Html a PDF", images[15], opcion16),
    ("Desbloquear PDF", images[16], opcion17),
    ("Proteger PDF", images[17], opcion18),
    ("Ordenar PDF", images[18], opcion19),
    ("PDF a PDF/a", images[19], opcion20),
    ("Reparar PDF", images[20], opcion21),
    ("Enumerar páginas", images[21], opcion22),
    ("Escanear a PDF", images[22], opcion23),
    ("OCR PDF", images[23], opcion24),
    ("Comparar PDF", images[24], opcion25),
]

# Añadir botón para cambiar el modo
toggle_button = tk.Button(root, image=night_icon, command=toggle_mode, bg="white", fg="black")
toggle_button.pack(side="left", padx=10, pady=10, anchor='se')

# Añadir botones a la cuadrícula
max_columns = 5
for i, (text, image, command) in enumerate(opciones):
    row = i // max_columns
    column = i % max_columns
    button = create_button(text, image, command)
    button.grid(row=row, column=column, padx=10, pady=10)
    buttons.append(button)

# Cargar el icono principal de la ventana
icono_img = Image.open("img/icono.png")
icono_tk = ImageTk.PhotoImage(icono_img)
root.iconphoto(True, icono_tk)

# Label para mostrar los mensajes
message_label = tk.Label(root)
message_label.pack(pady=10)

# Etiqueta para mostrar el nombre del autor
author_label = tk.Label(root, text="AsierBajo", font=("Helvetica", 12, "bold"))
author_label.place(relx=1.0, rely=1.0, anchor="se", x=-10, y=-10)

root.mainloop()
